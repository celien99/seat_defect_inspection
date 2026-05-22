"""PPTX report generation for benchmark evaluation."""

from __future__ import annotations

import io
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from .config import BenchmarkConfig
from .metrics import (
    compute_binary_metrics,
    compute_confusion_matrix,
    compute_score_distributions,
    compute_timing_stats,
)
from .plots import (
    fig_to_base64,
    plot_confusion_matrix,
    plot_defect_type_metrics,
    plot_per_camera_metrics,
    plot_pr_curve,
    plot_roc_curve,
    plot_score_distribution,
)
from .schemas import BenchmarkSummary

_pptx = None


def _get_pptx():
    global _pptx
    if _pptx is None:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            _pptx = True
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX report generation. "
                "Install it with: pip install python-pptx"
            )
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN


def _add_title_slide(prs, title: str, subtitle: str) -> None:
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _get_pptx()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_metric_slide(prs, title: str, metrics: Dict[str, str]) -> None:
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _get_pptx()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_body = slide.placeholders[1].text_frame
    text_body.clear()
    for label, value in metrics.items():
        p = text_body.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = Pt(18)


def _add_chart_slide(prs, title: str, chart_b64: str) -> None:
    if not chart_b64:
        return
    import base64
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _get_pptx()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    img_data = base64.b64decode(chart_b64)
    img_stream = io.BytesIO(img_data)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5.5))


def _add_table_slide(prs, title: str, headers: List[str], rows: List[List[str]]) -> None:
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _get_pptx()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    if not rows:
        return

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * n_rows)).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)


def generate_pptx_report(
    summary: BenchmarkSummary,
    config: BenchmarkConfig,
    output_path: str,
) -> str:
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _get_pptx()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    all_labeled = [r for rd in summary.rounds for r in rd.records if r.sample.ground_truth_label is not None]
    now_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # Slide 1: Title
    _add_title_slide(prs, "Seat Defect Inspection\nBenchmark Report", f"Generated: {now_str}")

    # Slide 2: Executive Summary
    cm = compute_confusion_matrix(all_labeled) if all_labeled else None
    bm = compute_binary_metrics(cm) if cm and cm.total > 0 else None
    exec_metrics = {}
    if bm:
        exec_metrics = {
            "Accuracy (准确率)": f"{bm.accuracy * 100:.1f}%",
            "Precision (精准率)": f"{bm.precision * 100:.1f}%",
            "Recall (召回率)": f"{bm.recall * 100:.1f}%",
            "F1 Score": f"{bm.f1 * 100:.1f}%",
            "Miss Rate (漏检率)": f"{bm.miss_rate * 100:.1f}%",
            "False Alarm (错检率)": f"{bm.false_alarm_rate * 100:.1f}%",
        }
    _add_metric_slide(prs, "Executive Summary (执行摘要)", exec_metrics)

    # Slide 3: Dataset Composition
    headers = ["Round", "Samples", "OK", "NG", "Ground Truth"]
    rows = [
        [rd.composition.round_name,
         str(rd.composition.sample_count),
         str(rd.composition.ok_count),
         str(rd.composition.ng_count),
         rd.composition.ground_truth_source]
        for rd in summary.rounds
    ]
    _add_table_slide(prs, "Dataset Composition (数据构成)", headers, rows)

    # Slide 4: Per-round metrics
    for rd in summary.rounds:
        if rd.confusion and rd.confusion.total > 0:
            chart_b64 = plot_confusion_matrix(rd.confusion, title=f"Confusion Matrix — {rd.round_name}")
            _add_chart_slide(prs, f"Confusion Matrix — {rd.round_name}", chart_b64)

    # Slide: Per-camera chart
    all_cam_metrics = []
    for rd in summary.rounds:
        all_cam_metrics.extend(rd.per_camera)
    if all_cam_metrics:
        # Deduplicate by camera_id
        seen_cams = set()
        unique_metrics = []
        for m in all_cam_metrics:
            if m.camera_id not in seen_cams:
                seen_cams.add(m.camera_id)
                unique_metrics.append(m)
        cam_chart = plot_per_camera_metrics(unique_metrics, title="Per-Camera Metrics")
        _add_chart_slide(prs, "Per-Camera Analysis (按机位分析)", cam_chart)

    # Slide: Score distribution
    score_dists = compute_score_distributions(all_labeled) if all_labeled else []
    if score_dists:
        s_chart = plot_score_distribution(score_dists)
        _add_chart_slide(prs, "Score Distribution (分值分布)", s_chart)

    # Slide: ROC/PR
    for rd in summary.rounds:
        if rd.roc_curve and rd.roc_curve.points:
            roc_chart = plot_roc_curve(rd.roc_curve, title=f"ROC — {rd.round_name}")
            pr_chart = plot_pr_curve(rd.pr_curve, title=f"PR — {rd.round_name}")
            _add_chart_slide(prs, f"ROC/PR Curves — {rd.round_name}", roc_chart)

    # Slide: Failure cases summary
    failure_cases = []
    for rd in summary.rounds:
        failure_cases.extend(rd.failure_cases)
    if failure_cases:
        fc_headers = ["Part ID", "GT", "Predicted", "Type", "Reason"]
        fc_rows = []
        for fc in failure_cases[:15]:
            fc_rows.append([
                fc.sample.part_id,
                fc.sample.ground_truth_label or "?",
                fc.predicted_status,
                fc.sample.ground_truth_defect_type or "N/A",
                fc.decision_reason[:60],
            ])
        _add_table_slide(prs, f"Failure Cases ({len(failure_cases)} total)", fc_headers, fc_rows)

    # Slide: Recommendations
    rec_items = {}
    if bm:
        if bm.miss_rate > 0.05:
            rec_items["High Miss Rate"] = f"{bm.miss_rate * 100:.1f}% — Review decision thresholds"
        if bm.false_alarm_rate > 0.1:
            rec_items["High False Alarm"] = f"{bm.false_alarm_rate * 100:.1f}% — Consider raising threshold"
        if not rec_items:
            rec_items["Status"] = "Pipeline performance is within acceptable ranges."
    _add_metric_slide(prs, "Recommendations (优化建议)", rec_items)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[benchmark] PPTX report saved to: {output_path}")
    return str(output_path)
